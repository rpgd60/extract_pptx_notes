from pptx import Presentation
import click

@click.command()
@click.argument('pres-file', type=str)
@click.option('--header', default='', help='Header to display.')
@click.option('--print-all/--no-print-all', '-p', default=False, help='If True, display title for slides without notes')
def extract_notes(header, pres_file, print_all):
    """Extract text from the slide notes in a power point presentation using python-pptx \n 
       Display also an optional informative header, if included in the options"""
    prs = Presentation(pres_file)
    # text_runs will be populated with a list of strings
    # one for each text run in presentation
    text_runs = []
    slide_num = 0
    num_notes = 0
    num_slides = len(prs.slides)
    if len(header) > 0:
        print(f'# {header}')
    for slide in prs.slides:
        slide_title = slide.shapes.title.text
        slide_num += 1
        slide_has_non_empty_note = False
        if slide.has_notes_slide:
            da_note = slide.notes_slide.notes_text_frame.text
            if not da_note.isspace() and len(da_note) > 0:
                slide_has_non_empty_note = True
                num_notes += 1
        if print_all | slide_has_non_empty_note:
            print(f'## [{slide_num}/{num_slides}] - {slide_title}\n')
            if slide_has_non_empty_note:
                print(f'### Note:\n{da_note.strip()}\n')
    print(f'{header}: Found {num_notes} slides with non-empty notes in {pres_file}')

if __name__ == "__main__":
    extract_notes()